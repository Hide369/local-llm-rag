import io

import pytest
from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from ingest.parsers import UnsupportedFormatError, parse
from ingest.parsers.docx_parser import parse_docx
from ingest.parsers.pptx_parser import parse_pptx


@pytest.fixture
def docx_path(tmp_path):
    """テスト用のdocxをその場で生成する（バイナリをリポジトリに置かないため）。"""
    doc = Document()
    doc.add_paragraph("会議名：キックオフ")
    doc.add_paragraph("")  # 空段落は無視されること
    doc.add_paragraph("決定事項：RAGを導入する")
    path = tmp_path / "議事録.docx"
    doc.save(path)
    return path


def _png_bytes(width, height, color="green"):
    buf = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture
def pptx_path(tmp_path):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank)
    box = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "一枚目のタイトル"

    slide2 = prs.slides.add_slide(blank)
    box2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box2.text_frame.text = "二枚目の本文"
    slide2.notes_slide.notes_text_frame.text = "発表者ノート"

    path = tmp_path / "資料.pptx"
    prs.save(path)
    return path


def test_docx_becomes_a_single_document_unit(docx_path):
    units = parse_docx(docx_path)
    assert len(units) == 1
    assert units[0].location_type == "document"
    assert units[0].location == 0


def test_docx_joins_paragraphs_and_drops_empty_ones(docx_path):
    text = parse_docx(docx_path)[0].text
    assert "会議名：キックオフ" in text
    assert "決定事項：RAGを導入する" in text
    assert "\n\n" not in text


def test_pptx_produces_one_unit_per_slide(pptx_path):
    units = parse_pptx(pptx_path)
    assert len(units) == 2
    assert [u.location for u in units] == [1, 2]
    assert all(u.location_type == "slide" for u in units)


def test_pptx_includes_speaker_notes(pptx_path):
    """ノート欄には本文に書かれない補足が入るため取り込み対象にする。"""
    assert "発表者ノート" in parse_pptx(pptx_path)[1].text


def test_pptx_slide_text_is_captured(pptx_path):
    assert "一枚目のタイトル" in parse_pptx(pptx_path)[0].text


def test_office_parsers_are_not_marked_as_ocr(docx_path, pptx_path):
    assert parse_docx(docx_path)[0].ocr is False
    assert all(u.ocr is False for u in parse_pptx(pptx_path))


@pytest.fixture
def rich_pptx_path(tmp_path):
    """実データのスライド11を模した1枚。

    読み順・ノイズ除去・タイトル複写・分割をまとめて確かめる。シェイプは
    わざとXML順と視覚順がずれるように置く（実データも同じ形だった）。
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def box(text, top_inches):
        shape = slide.shapes.add_textbox(
            Inches(1), Inches(top_inches), Inches(6), Inches(1)
        )
        shape.text_frame.text = text
        return shape

    box("©2026 OpenUp Next Engineer Inc.", 6.0)
    box("1", 5.5)
    box("あ" * 150, 2.0)
    box("1-6. 生成AI活用のポイント", 0.5)
    box("い" * 120, 3.0)

    path = tmp_path / "セミナー.pptx"
    prs.save(path)
    return path


@pytest.fixture
def pptx_with_large_picture(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(4), Inches(1))
    box.text_frame.text = "タイトル"
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes(800, 600)), Inches(1), Inches(2), Inches(4), Inches(3)
    )
    path = tmp_path / "図あり.pptx"
    prs.save(path)
    return path


@pytest.fixture
def pptx_with_small_picture(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(4), Inches(1))
    box.text_frame.text = "タイトル"
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes(40, 40, "blue")), Inches(0.1), Inches(0.1), Inches(0.3), Inches(0.3)
    )
    path = tmp_path / "小さい画像.pptx"
    prs.save(path)
    return path


def test_shapes_are_read_in_visual_order_not_xml_order(rich_pptx_path):
    """XML順ではタイトルが4番目に来る。読み順に直さないとタイトルを取り違える。"""
    units = parse_pptx(rich_pptx_path)
    assert units[0].text.startswith("1-6. 生成AI活用のポイント")


def test_the_slide_is_split_into_several_units(rich_pptx_path):
    """150字と120字は合わせて目標200字を超えるので別グループになる。"""
    units = parse_pptx(rich_pptx_path)
    assert len(units) == 2
    assert "あ" in units[0].text and "い" not in units[0].text
    assert "い" in units[1].text and "あ" not in units[1].text


def test_the_title_is_copied_into_every_unit(rich_pptx_path):
    """分割後のチャンクが単独で何の話か分かるようにする。"""
    units = parse_pptx(rich_pptx_path)
    assert all(u.text.startswith("1-6. 生成AI活用のポイント\n") for u in units)


def test_a_shape_is_never_split_across_units(rich_pptx_path):
    """シェイプは作成者が区切った意味のまとまり。跨いで切ると文脈が壊れる。"""
    units = parse_pptx(rich_pptx_path)
    assert units[0].text.count("あ") == 150
    assert units[1].text.count("い") == 120


def test_copyright_and_page_number_lines_are_dropped(rich_pptx_path):
    """全43枚に同じフッタが入っており、ベクトルを一様に濁らせる。"""
    joined = "\n".join(u.text for u in parse_pptx(rich_pptx_path))
    assert "©" not in joined
    assert "OpenUp" not in joined
    assert "\n1\n" not in joined and not joined.endswith("\n1")


def test_every_unit_of_a_slide_keeps_the_slide_number(rich_pptx_path):
    """出典が「スライド11」であることは分割後も変わらない。"""
    units = parse_pptx(rich_pptx_path)
    assert all(u.location == 1 and u.location_type == "slide" for u in units)


def test_a_shape_larger_than_the_target_becomes_its_own_unit(tmp_path):
    """実データには538字のシェイプがある。200字は目標であって上限ではない。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for text, top in (("見出しの行", 0.5), ("う" * 500, 2.0)):
        shape = slide.shapes.add_textbox(Inches(1), Inches(top), Inches(6), Inches(1))
        shape.text_frame.text = text
    path = tmp_path / "長い.pptx"
    prs.save(path)

    units = parse_pptx(path)
    assert len(units) == 1
    assert units[0].text.count("う") == 500


def test_text_inside_grouped_shapes_is_captured(tmp_path):
    """グループ内のテキストは現行の実装が無言で捨てていた（spec 7.7）。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(1))
    title.text_frame.text = "手順の説明"
    group = slide.shapes.add_group_shape()
    inner = group.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    inner.text_frame.text = "Step 1 レビュー依頼"
    path = tmp_path / "グループ.pptx"
    prs.save(path)

    assert "Step 1 レビュー依頼" in parse_pptx(path)[0].text


def test_a_slide_with_only_a_title_still_produces_a_unit(tmp_path):
    """本文が空でもユニットを落とさない（spec 7.6）。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "章の扉ページ"
    path = tmp_path / "扉.pptx"
    prs.save(path)

    units = parse_pptx(path)
    assert len(units) == 1
    assert units[0].text == "章の扉ページ"


def test_remaining_lines_of_the_title_shape_are_kept_as_body(tmp_path):
    """タイトルと副題が同じテキストボックスに入っている場合、副題を捨てない。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(2))
    box.text_frame.text = "生成AI活用セミナー\nAIエージェントとAI駆動開発を学ぼう"
    path = tmp_path / "扉2.pptx"
    prs.save(path)

    units = parse_pptx(path)
    assert len(units) == 1
    assert units[0].text == "生成AI活用セミナー\nAIエージェントとAI駆動開発を学ぼう"


def test_dispatch_routes_by_suffix(docx_path, pptx_path):
    assert parse(docx_path)[0].location_type == "document"
    assert parse(pptx_path)[0].location_type == "slide"


def test_dispatch_rejects_unsupported_suffix(tmp_path):
    other = tmp_path / "memo.txt"
    other.write_text("本文", encoding="utf-8")
    with pytest.raises(UnsupportedFormatError):
        parse(other)


def test_caption_image_not_called_when_not_provided(pptx_with_large_picture):
    def _fail(_bytes):
        raise AssertionError("caption_imageが未指定なら呼ばれてはいけない")

    units = parse_pptx(pptx_with_large_picture, caption_image=None)
    assert all("[図の説明]" not in u.text for u in units)
    assert all(u.vlm is False for u in units)


def test_large_picture_is_captioned(pptx_with_large_picture):
    units = parse_pptx(
        pptx_with_large_picture, caption_image=lambda _bytes: "緑色の図表です。"
    )
    assert any("[図の説明] 緑色の図表です。" in u.text for u in units)
    assert any(u.vlm for u in units)


def test_small_picture_is_not_captioned(pptx_with_small_picture):
    def _fail(_bytes):
        raise AssertionError("閾値未満の画像でcaption_imageが呼ばれてはいけない")

    units = parse_pptx(pptx_with_small_picture, caption_image=_fail)
    assert all("[図の説明]" not in u.text for u in units)
    assert all(u.vlm is False for u in units)


def test_picture_caption_failure_is_skipped_with_a_warning(pptx_with_large_picture, capsys):
    from ingest.vlm import VlmError

    def _raise(_bytes):
        raise VlmError("boom")

    units = parse_pptx(pptx_with_large_picture, caption_image=_raise)
    assert all("[図の説明]" not in u.text for u in units)
    assert "boom" in capsys.readouterr().err
    assert any("タイトル" in u.text for u in units)
