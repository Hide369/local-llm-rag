"""PowerPointのテキスト抽出。

スライドは話題の単位としては粗すぎる。実測では1スライドに7つの話題が同居し、
チャンクのベクトルがスライド全体の平均に薄まって、その中の1話題を指す質問で
順位を落としていた。同じチャンクが、スライドの表題そのままの質問には距離0.303で
1位、その中の1話題を指す質問には0.514で4位になる（spec 2.4）。

テキストボックスは作成者が視覚的に区切った意味のまとまりなので、これを分割の
単位にする。実測では43スライドが80ユニット（min 11 / max 547 / avg 156字）になる。
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from ingest.models import SLIDE, ParsedUnit

# 同じ行に並ぶ図表の要素を左から右へ並べるための丸め幅。0.25インチをEMUで表した値。
ROW_TOLERANCE = 228600

# ブロックを積んでいく目標文字数。上限ではない。単独でこれを超えるシェイプ
# （実データの最大は538字）は、それ1つで1グループになる。
GROUP_TARGET_CHARS = 200

# ロゴ・アイコン等の装飾画像を除外するための閾値（配置サイズ、インチ角）。
# PDFの画素サイズと違い、スライド上に配置された大きさで判定できる。実データでの
# 実測は design docの6節を参照。未実測の仮値であり、後日調整する前提。
MIN_PICTURE_INCHES = 1.0

_CAPTION_PREFIX = "[図の説明] "


def _walk(shapes):
    """グループ化されたシェイプの中まで辿る。

    トップレベルだけを見ると、グループ内のテキストが例外も警告もなく捨てられる。
    実データではスライド26の「Step 1」〜「Step 4」が該当した。
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk(shape.shapes)
        else:
            yield shape


def _position(shape):
    """読み順（行優先）で並べるための整列キー。

    XML上の並び順は視覚的な読み順と一致しない。実データではタイトルの次に
    ページ番号が来る。top を丸めて行にまとめ、その中を left で並べる。
    """
    return (round((shape.top or 0) / ROW_TOLERANCE), shape.left or 0)


def _is_captionable_picture(shape) -> bool:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    return (
        Emu(shape.width).inches >= MIN_PICTURE_INCHES
        and Emu(shape.height).inches >= MIN_PICTURE_INCHES
    )


def _picture_block(shape, slide_number: int, path_name: str, caption_image) -> str | None:
    try:
        caption = caption_image(shape.image.blob)
    except Exception as error:
        print(
            f"警告: 画像の説明取得に失敗しました（{path_name} スライド{slide_number}）: {error}",
            file=sys.stderr,
        )
        return None
    if not caption.strip() or caption.strip() == "装飾画像":
        return None
    return f"{_CAPTION_PREFIX}{caption}"


def _split_title(blocks: list[str]) -> tuple[str, list[str]]:
    """先頭のテキストブロックをタイトルとして取り出す。

    画像キャプションが先頭に来ても、それをスライドのタイトルとして全ユニットへ
    複写してしまわないようにする（上部に図を置くレイアウトとの相互作用）。
    """
    for index, block in enumerate(blocks):
        if not block.startswith(_CAPTION_PREFIX):
            title, _, remainder = block.partition("\n")
            body = blocks[:index] + ([remainder] if remainder else []) + blocks[index + 1 :]
            return title, body
    # 全ブロックが画像キャプションだった場合は先頭をそのままタイトルにする。
    title, _, remainder = blocks[0].partition("\n")
    return title, ([remainder] if remainder else []) + blocks[1:]


def _clean(text: str, slide_number: int) -> str:
    """内容を持たない行を落とす。

    実データでは全43枚にコピーライトのフッタが、38枚にページ番号だけの行が
    入っていた。全スライドに同じ文字列が入るとベクトルが一様に濁る。BM25側は
    IDFが押さえるが、ベクトル側には効かないため取り込み時に落とす。
    ページ番号は location メタデータが持っているので失われない。
    """
    lines = [
        line
        for line in (raw.strip() for raw in text.split("\n"))
        if line and not line.startswith("©") and line != str(slide_number)
    ]
    return "\n".join(lines)


def _blocks(slide, slide_number: int, path_name: str = "", caption_image=None) -> list[str]:
    """読み順に並べた、内容のあるシェイプのテキスト。"""
    content_shapes = [
        shape
        for shape in _walk(slide.shapes)
        if shape.has_text_frame
        or (caption_image is not None and _is_captionable_picture(shape))
    ]
    content_shapes.sort(key=_position)

    blocks = []
    for shape in content_shapes:
        if shape.has_text_frame:
            cleaned = _clean(shape.text_frame.text, slide_number)
            if cleaned:
                blocks.append(cleaned)
        else:
            block = _picture_block(shape, slide_number, path_name, caption_image)
            if block is not None:
                blocks.append(block)

    # ノート欄には本文に書かれていない補足や発表意図が入るため取り込む。
    # 位置情報を持たないので読み順の最後尾に置く。
    if slide.has_notes_slide:
        notes = _clean(slide.notes_slide.notes_text_frame.text, slide_number)
        if notes:
            blocks.append(notes)
    return blocks


def _group(blocks: list[str]) -> list[str]:
    """目標文字数を目安にブロックをまとめる。シェイプは分割しない。"""
    groups: list[str] = []
    current: list[str] = []
    length = 0
    for block in blocks:
        # 足すと目標を超えるなら、足す前に閉じる。ただし空のときは閉じない。
        # 単独で目標を超えるシェイプが、それ1つで1グループになるようにするため。
        if current and length + len(block) > GROUP_TARGET_CHARS:
            groups.append("\n".join(current))
            current, length = [], 0
        current.append(block)
        length += len(block)
    if current:
        groups.append("\n".join(current))
    return groups


def parse_pptx(path: Path, caption_image=None) -> list[ParsedUnit]:
    units: list[ParsedUnit] = []
    for number, slide in enumerate(Presentation(path).slides, start=1):
        blocks = _blocks(slide, number, path_name=path.name, caption_image=caption_image)
        if not blocks:
            continue
        # 先頭シェイプの1行目をタイトルとする。この資料では slide.shapes.title が
        # 43枚すべて None だった（タイトルプレースホルダを使っていない）。
        # 残りの行は捨てず、最初の本文ブロックとして扱う。
        title, body = _split_title(blocks)
        # タイトルを全ユニットへ複写する。分割後のチャンクが単独で何の話か
        # 分かるようにするため。「事前学習済モデルに追加学習させ、LLMを再生成する」
        # だけを見ても、何の定義か分からない。
        # 本文が空のスライド（章の扉）はタイトルだけで1ユニットにする。
        for group in _group(body) or [""]:
            text = f"{title}\n{group}" if group else title
            units.append(
                ParsedUnit(
                    text=text,
                    location_type=SLIDE,
                    location=number,
                    vlm=_CAPTION_PREFIX in text,
                )
            )
    return units
